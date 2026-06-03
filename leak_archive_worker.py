#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
leak_archive_worker.py

Pipeline:
download -> ClamAV scan archive -> extract -> ClamAV scan extracted folder
-> inventory/score files -> import useful text-like files into crawler.db
"""

import logging
import csv
import json
import argparse
import hashlib
import mimetypes
import os
import re
import shutil
import sqlite3
import subprocess
import time
import zipfile
from datetime import datetime
from pathlib import Path
from urllib.request import Request, urlopen
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse

logging.getLogger("pypdf").setLevel(logging.ERROR)
logging.getLogger("PyPDF2").setLevel(logging.ERROR)


# --- nested archive support patch ---
NESTED_ARCHIVE_EXTS = {".zip", ".7z", ".rar", ".tar", ".gz", ".tgz"}

def is_macos_resource_fork(path):
    import os
    parts = path.replace("\\", "/").split("/")
    base = os.path.basename(path)
    return "__MACOSX" in parts or base.startswith("._") or base == ".DS_Store"

def expand_nested_archives(root_dir, max_depth=3):
    import os, zipfile, subprocess, shutil
    from pathlib import Path

    root = Path(root_dir)
    expanded = 0

    for depth in range(max_depth):
        found = False

        for f in list(root.rglob("*")):
            if not f.is_file():
                continue

            if is_macos_resource_fork(str(f)):
                continue

            ext = f.suffix.lower()
            if ext not in NESTED_ARCHIVE_EXTS:
                continue

            out_dir = f.with_suffix("")
            out_dir = f.parent / (out_dir.name + "__nested_extract")

            if out_dir.exists():
                continue

            out_dir.mkdir(parents=True, exist_ok=True)

            try:
                if ext == ".zip":
                    with zipfile.ZipFile(f, "r") as z:
                        z.extractall(out_dir)
                else:
                    shutil.which("7z") or (_ for _ in ()).throw(RuntimeError("7z not installed"))
                    subprocess.run(["7z", "x", "-y", f"-o{out_dir}", str(f)], check=True)

                found = True
                expanded += 1
                print(f"[nested-archive] extracted {f} -> {out_dir}")

            except Exception as e:
                print(f"[nested-archive] failed {f}: {e}")

        if not found:
            break

    return expanded
# --- end nested archive support patch ---


BASE_DIR = Path(__file__).resolve().parent
DB_PATH = BASE_DIR / "crawler.db"
BASE_TMP = Path("/tmp/leak_processing")

ALLOWED_EXTS = {
    ".csv", ".tsv", ".txt", ".json", ".ndjson", ".sql", ".log", ".xml",
    ".pdf", ".pptx"
}

DOCUMENT_EXTRACT_EXTS = {
    ".pdf", ".pptx"
}

SKIP_EXTS = {
    ".exe", ".dll", ".so", ".bin", ".dat", ".jpg", ".jpeg", ".png",
    ".gif", ".mp4", ".mp3", ".iso", ".img", ".bat", ".ps1",
    ".vbs", ".js", ".jar", ".scr", ".msi", ".com"
}

MAX_DOWNLOAD_BYTES = None
MAX_SAMPLE_BYTES = 1024 * 1024
MAX_FULL_FILE_BYTES = 500 * 1024 * 1024
MAX_EXTRACTED_FILES = 50000

EMAIL_RE = re.compile(r"\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b")
SSN_RE = re.compile(r"\b\d{3}-\d{2}-\d{4}\b")
PHONE_RE = re.compile(r"\b(?:\+?1[\s.\-]?)?\(?\d{3}\)?[\s.\-]?\d{3}[\s.\-]?\d{4}\b")
URL_RE = re.compile(r"https?://[^\s\"'<>]+")
DOMAIN_RE = re.compile(r"\b(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,}\b")


def utc_now():
    return datetime.utcnow().isoformat(timespec="seconds")


def connect_db():
    con = sqlite3.connect(str(DB_PATH), timeout=60)
    con.row_factory = sqlite3.Row
    con.execute("PRAGMA journal_mode=WAL")
    con.execute("PRAGMA synchronous=NORMAL")
    return con


def column_exists(con, table, column):
    try:
        return column in {r[1] for r in con.execute(f"PRAGMA table_info({table})").fetchall()}
    except Exception:
        return False


def add_column_if_missing(con, table, column, definition):
    if not column_exists(con, table, column):
        con.execute(f"ALTER TABLE {table} ADD COLUMN {column} {definition}")


def init_tables(con):
    con.executescript("""
    CREATE TABLE IF NOT EXISTS leak_archives (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        download_url TEXT,
        victim_name TEXT,
        actor TEXT,
        leak_name TEXT,
        source_type TEXT,
        local_path TEXT,
        sha256 TEXT,
        status TEXT DEFAULT 'queued',
        malware_detected INTEGER DEFAULT 0,
        malware_signature TEXT,
        scan_result TEXT,
        error TEXT,
        files_found INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0,
        entities_imported INTEGER DEFAULT 0,
        created_at TEXT,
        updated_at TEXT
    );

    CREATE TABLE IF NOT EXISTS leak_archive_files (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        archive_id INTEGER,
        file_path TEXT,
        file_ext TEXT,
        size_bytes INTEGER,
        detected_type TEXT,
        score INTEGER DEFAULT 0,
        status TEXT DEFAULT 'queued',
        reason TEXT,
        created_at TEXT
    );

    CREATE TABLE IF NOT EXISTS exposure_entities (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        value_plain TEXT,
        entity_type TEXT,
        times_seen INTEGER DEFAULT 1,
        first_seen TEXT,
        last_seen TEXT,
        UNIQUE(value_plain, entity_type)
    );

    CREATE TABLE IF NOT EXISTS exposure_occurrences (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        victim_name TEXT,
        actor TEXT,
        leak_name TEXT,
        entity_id INTEGER,
        value_plain TEXT,
        entity_type TEXT,
        source_file TEXT,
        archive_id INTEGER,
        created_at TEXT
    );

    CREATE INDEX IF NOT EXISTS idx_exposure_entities_value ON exposure_entities(value_plain);
    CREATE INDEX IF NOT EXISTS idx_exposure_entities_type ON exposure_entities(entity_type);
    CREATE INDEX IF NOT EXISTS idx_exposure_occ_actor ON exposure_occurrences(actor);
    CREATE INDEX IF NOT EXISTS idx_exposure_occ_victim ON exposure_occurrences(victim_name);
    CREATE INDEX IF NOT EXISTS idx_leak_archive_files_archive ON leak_archive_files(archive_id, score DESC);
    """)

    # Safe migrations if older tables already exist.
    for col, definition in [
        ("value_plain", "TEXT"),
        ("entity_type", "TEXT"),
        ("times_seen", "INTEGER DEFAULT 1"),
        ("first_seen", "TEXT"),
        ("last_seen", "TEXT"),
    ]:
        add_column_if_missing(con, "exposure_entities", col, definition)

    for col, definition in [
        ("victim_name", "TEXT"),
        ("actor", "TEXT"),
        ("leak_name", "TEXT"),
        ("entity_id", "INTEGER"),
        ("value_plain", "TEXT"),
        ("entity_type", "TEXT"),
        ("source_file", "TEXT"),
        ("archive_id", "INTEGER"),
        ("archive_file_id", "INTEGER"),
        ("row_number", "INTEGER"),
        ("context", "TEXT"),
        ("raw_column", "TEXT"),
        ("created_at", "TEXT"),
    ]:
        add_column_if_missing(con, "exposure_occurrences", col, definition)

    for col, definition in [
        ("rows_seen", "INTEGER DEFAULT 0"),
        ("error", "TEXT"),
        ("inner_path", "TEXT"),
        ("inner_filename", "TEXT"),
        ("parent_folder", "TEXT"),
        ("extension", "TEXT"),
        ("compressed_size", "INTEGER DEFAULT 0"),
        ("uncompressed_size", "INTEGER DEFAULT 0"),
        ("processable", "INTEGER DEFAULT 0"),
        ("processed", "INTEGER DEFAULT 0"),
    ]:
        add_column_if_missing(con, "leak_archive_files", col, definition)

    con.commit()


def update_archive(con, archive_id, **fields):
    fields["updated_at"] = utc_now()
    keys = list(fields.keys())
    vals = [fields[k] for k in keys] + [archive_id]
    con.execute(
        f"UPDATE leak_archives SET {', '.join([f'{k}=?' for k in keys])} WHERE id=?",
        vals
    )
    con.commit()


def get_archive(con, archive_id):
    return con.execute("SELECT * FROM leak_archives WHERE id=?", (archive_id,)).fetchone()


def sha256_file(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def download_file(url, dest):
    req = Request(url, headers={"User-Agent": "Mozilla/5.0"})
    total = 0

    with urlopen(req, timeout=90) as r, open(dest, "wb") as f:
        while True:
            chunk = r.read(1024 * 1024)
            if not chunk:
                break
            total += len(chunk)
            if MAX_DOWNLOAD_BYTES is not None and total > MAX_DOWNLOAD_BYTES:
                raise RuntimeError("Download exceeded MAX_DOWNLOAD_BYTES")
            f.write(chunk)


def clam_scan(path):
    commands = [
        ["clamdscan", "--fdpass", str(path)],
        ["clamscan", "-r", str(path)],
    ]

    last_output = ""
    last_code = -1

    for cmd in commands:
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=900
            )
            output = (result.stdout or "") + (result.stderr or "")
            last_output = output
            last_code = result.returncode

            infected = "FOUND" in output
            sig = None
            m = re.search(r": (.+?) FOUND", output)
            if m:
                sig = m.group(1)

            # ClamAV return codes: 0 clean, 1 infected, 2 error
            if result.returncode in (0, 1):
                return {
                    "clean": result.returncode == 0 and not infected,
                    "infected": infected,
                    "signature": sig,
                    "output": output,
                    "return_code": result.returncode,
                    "scanner": cmd[0],
                }

        except FileNotFoundError:
            continue
        except Exception as e:
            last_output = str(e)
            last_code = -1

    return {
        "clean": False,
        "infected": False,
        "signature": None,
        "output": last_output or "No ClamAV scanner available",
        "return_code": last_code,
        "scanner": "none",
    }


def safe_zip_extract(zip_path, extract_dir):
    extract_dir = extract_dir.resolve()
    count = 0

    with zipfile.ZipFile(zip_path, "r") as z:
        for member in z.infolist():
            count += 1
            if count > MAX_EXTRACTED_FILES:
                raise RuntimeError("Archive exceeded MAX_EXTRACTED_FILES")

            target = (extract_dir / member.filename).resolve()
            if not str(target).startswith(str(extract_dir)):
                raise RuntimeError(f"Blocked unsafe archive path: {member.filename}")

            if member.is_dir():
                target.mkdir(parents=True, exist_ok=True)
                continue

            target.parent.mkdir(parents=True, exist_ok=True)
            with z.open(member) as src, open(target, "wb") as dst:
                shutil.copyfileobj(src, dst)


def extract_archive(archive_path, extract_dir):
    extract_dir.mkdir(parents=True, exist_ok=True)
    suffix = archive_path.suffix.lower()

    if suffix == ".zip":
        safe_zip_extract(archive_path, extract_dir)
        return

    if suffix == ".7z":
        try:
            import py7zr
        except ImportError:
            raise RuntimeError("7z file detected, but py7zr is not installed. Install with: pip install py7zr")
        with py7zr.SevenZipFile(archive_path, mode="r") as z:
            z.extractall(path=extract_dir)
        return

    if suffix in ALLOWED_EXTS:
        shutil.copy2(archive_path, extract_dir / archive_path.name)
        return

    raise RuntimeError(f"Unsupported archive/file type: {suffix}")


def is_binary(path):
    try:
        with open(path, "rb") as f:
            sample = f.read(8192)
        return b"\x00" in sample
    except Exception:
        return True


def extract_pdf_text(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("PDF file detected, but pypdf is not installed. Install with: pip install pypdf")

    text = []
    reader = PdfReader(str(path))

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)

    return "\n".join(text)


def extract_pptx_text(path):
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("PPTX file detected, but python-pptx is not installed. Install with: pip install python-pptx")

    text = []
    prs = Presentation(str(path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)

    return "\n".join(text)


def extract_document_text(path):
    ext = path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf_text(path)

    if ext == ".pptx":
        return extract_pptx_text(path)

    raise RuntimeError(f"Unsupported document extraction type: {ext}")


def score_file(path):
    ext = path.suffix.lower()
    size = path.stat().st_size
    name = path.name.lower()
    score = 0
    reasons = []

    if ext in ALLOWED_EXTS:
        score += 20
        reasons.append("allowed_extension")

    if ext in SKIP_EXTS:
        score -= 100
        reasons.append("blocked_extension")

    if size <= 0:
        score -= 100
        reasons.append("empty_file")

    if size > MAX_FULL_FILE_BYTES:
        score -= 50
        reasons.append("too_large")

    if any(x in name for x in ["data", "users", "customers", "clients", "leak", "dump", "export", "database", "member"]):
        score += 10
        reasons.append("useful_filename")

    if any(x in name for x in ["readme", "license", "instructions", "about"]):
        score -= 20
        reasons.append("low_value_filename")

    if ext in DOCUMENT_EXTRACT_EXTS:
        score += 50
        reasons.append("extractable_document")

        detected = mimetypes.guess_type(str(path))[0] or "application/octet-stream"

        try:
            sample = extract_document_text(path)[:MAX_SAMPLE_BYTES]
        except Exception as e:
            score -= 80
            reasons.append(f"document_extract_error:{e}")
            return score, detected, ",".join(reasons)

        if not sample.strip():
            score -= 30
            reasons.append("no_extractable_text")

    else:
        if is_binary(path):
            score -= 100
            reasons.append("binary_file")
            return score, "binary", ",".join(reasons)

        try:
            sample = path.read_bytes()[:MAX_SAMPLE_BYTES].decode("utf-8", errors="ignore")
        except Exception:
            return -100, "unreadable", "decode_error"

    if EMAIL_RE.search(sample):
        score += 25
        reasons.append("email_found")
    if SSN_RE.search(sample):
        score += 25
        reasons.append("ssn_found")
    if PHONE_RE.search(sample):
        score += 10
        reasons.append("phone_found")
    if URL_RE.search(sample):
        score += 10
        reasons.append("url_found")

    lines = sample.splitlines()
    first_line = lines[0] if lines else ""
    if "," in first_line or "\t" in first_line or "|" in first_line:
        score += 10
        reasons.append("delimited_header")

    detected = mimetypes.guess_type(str(path))[0] or "text/plain"
    return score, detected, ",".join(reasons)


def inventory_files(con, archive_id, extract_dir):
    count = 0
    for path in extract_dir.rglob("*"):
        if not path.is_file():
            continue
        if is_macos_resource_fork(str(path)):
            continue
        count += 1
        if count > MAX_EXTRACTED_FILES:
            raise RuntimeError("Extracted folder exceeded MAX_EXTRACTED_FILES")

        score, detected, reason = score_file(path)

        con.execute("""
            INSERT INTO leak_archive_files
            (archive_id, file_path, file_ext, size_bytes, detected_type, score, status, reason)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            archive_id,
            str(path),
            path.suffix.lower(),
            path.stat().st_size,
            detected,
            score,
            "queued" if score >= 20 else "skipped",
            reason,
        ))

    con.commit()
    return count


def upsert_entity(con, value, entity_type):
    row = con.execute("""
        SELECT id FROM exposure_entities
        WHERE value_plain=? AND entity_type=?
    """, (value, entity_type)).fetchone()

    if row:
        con.execute("""
            UPDATE exposure_entities
            SET times_seen=COALESCE(times_seen,0)+1,
                last_seen=CURRENT_TIMESTAMP
            WHERE id=?
        """, (row["id"],))
        return row["id"]

    cur = con.execute("""
        INSERT INTO exposure_entities
        (value_plain, entity_type, times_seen, first_seen, last_seen)
        VALUES (?, ?, 1, CURRENT_TIMESTAMP, CURRENT_TIMESTAMP)
    """, (value, entity_type))
    return cur.lastrowid


def insert_occurrence(con, archive, entity_id, value, entity_type, source_file, archive_file_id=None, row_number=None, context=None, raw_column=None):
    if context is not None:
        context = str(context)
        if len(context) > 20000:
            context = context[:20000] + " ...[truncated]"

    con.execute("""
        INSERT INTO exposure_occurrences
        (victim_name, actor, leak_name, entity_id, value_plain, entity_type,
         source_file, archive_id, archive_file_id, row_number, context, raw_column, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, CURRENT_TIMESTAMP)
    """, (
        archive["victim_name"],
        archive["actor"],
        archive["leak_name"],
        entity_id,
        value,
        entity_type,
        source_file,
        archive["id"],
        archive_file_id,
        row_number,
        context,
        raw_column,
    ))


def extract_entities_from_text(text):
    found = []

    for value in EMAIL_RE.findall(text):
        found.append((value.lower(), "email"))
    for value in SSN_RE.findall(text):
        found.append((value, "ssn"))
    for value in PHONE_RE.findall(text):
        found.append((value, "phone"))
    for value in URL_RE.findall(text):
        found.append((value, "url"))
    for value in DOMAIN_RE.findall(text):
        if "@" not in value and not value.lower().startswith(("http.", "https.")):
            found.append((value.lower(), "domain"))

    return found


def process_file(con, archive, file_row):
    path = Path(file_row["file_path"])
    imported = 0

    try:
        if path.suffix.lower() in {".csv", ".tsv"}:
            delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
            with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
                reader = csv.DictReader(f, delimiter=delimiter)
                for row_number, row in enumerate(reader, start=2):
                    clean_row = {str(k): ("" if v is None else str(v)) for k, v in row.items() if k is not None}
                    context = json.dumps(clean_row, ensure_ascii=False)

                    for col_name, cell in clean_row.items():
                        if not cell:
                            continue

                        for value, entity_type in extract_entities_from_text(cell):
                            entity_id = upsert_entity(con, value, entity_type)
                            insert_occurrence(
                                con,
                                archive,
                                entity_id,
                                value,
                                entity_type,
                                str(path),
                                archive_file_id=file_row["id"],
                                row_number=row_number,
                                context=context,
                                raw_column=col_name,
                            )
                            imported += 1

                    if imported and imported % 5000 == 0:
                        con.commit()
        else:
            if path.suffix.lower() in DOCUMENT_EXTRACT_EXTS:
                lines = extract_document_text(path).splitlines()
            else:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    lines = f.readlines()

            for row_number, line in enumerate(lines, start=1):
                context = line.strip()
                for value, entity_type in extract_entities_from_text(line):
                    entity_id = upsert_entity(con, value, entity_type)
                    insert_occurrence(
                        con,
                        archive,
                        entity_id,
                        value,
                        entity_type,
                        str(path),
                        archive_file_id=file_row["id"],
                        row_number=row_number,
                        context=context,
                        raw_column=value,
                    )
                    imported += 1

            if imported and imported % 5000 == 0:
                con.commit()

        con.execute("""
            UPDATE leak_archive_files
            SET status='processed', reason=?, rows_seen=COALESCE(rows_seen,0)+?
            WHERE id=?
        """, (f"imported={imported}", imported, file_row["id"]))
        con.commit()
        return imported

    except Exception as e:
        con.execute("""
            UPDATE leak_archive_files
            SET status='error', reason=?
            WHERE id=?
        """, (str(e), file_row["id"]))
        con.commit()
        return 0


def process_archive(archive_id):
    con = connect_db()
    init_tables(con)

    archive = get_archive(con, archive_id)
    if not archive:
        raise RuntimeError(f"Archive job not found: {archive_id}")

    work_dir = BASE_TMP / f"archive_{archive_id}_{int(time.time())}"
    download_dir = work_dir / "download"
    extract_dir = work_dir / "extract"
    download_dir.mkdir(parents=True, exist_ok=True)

    try:
        update_archive(con, archive_id, status="downloading", error=None)

        url_path = (archive["download_url"] or "").split("?")[0]
        ext = Path(url_path).suffix.lower()
        filename = f"download_{archive_id}{ext if ext else ''}"
        archive_path = download_dir / filename

        if (archive["download_url"] or "").startswith("local://"):
            local_name = archive["download_url"].replace("local://", "", 1)
            candidates = []
            if "local_path" in archive.keys() and archive["local_path"]:
                candidates.append(Path(archive["local_path"]))
            candidates.extend([
                BASE_DIR / "stealer_uploads" / local_name,
                BASE_DIR / "manual_uploads" / local_name,
                BASE_DIR / local_name,
            ])
            local_src = next((c for c in candidates if c.exists()), None)

            if not local_src:
                tried = ", ".join(str(c) for c in candidates)
                raise RuntimeError(f"Local upload file not found. Tried: {tried}")

            shutil.copy2(local_src, archive_path)
        else:
            download_file(resolve_download_url(archive["download_url"]), archive_path)

        update_archive(con, archive_id, local_path=str(archive_path), sha256=sha256_file(archive_path))

        update_archive(con, archive_id, status="scanning_archive")
        scan = clam_scan(archive_path)

        if scan["infected"]:
            update_archive(
                con,
                archive_id,
                status="infected",
                malware_detected=1,
                malware_signature=scan["signature"],
                scan_result=scan["output"],
            )
            return

        if not scan["clean"]:
            update_archive(
                con,
                archive_id,
                status="scan_error",
                scan_result=scan["output"],
                error=f"ClamAV scan failed using {scan.get('scanner')}",
            )
            return

        update_archive(con, archive_id, status="extracting")
        extract_archive(archive_path, extract_dir)
        expand_nested_archives(extract_dir)

        update_archive(con, archive_id, status="scanning_extracted")
        scan2 = clam_scan(extract_dir)

        if scan2["infected"]:
            update_archive(
                con,
                archive_id,
                status="infected_extracted",
                malware_detected=1,
                malware_signature=scan2["signature"],
                scan_result=scan2["output"],
            )
            return

        if not scan2["clean"]:
            update_archive(
                con,
                archive_id,
                status="scan_error",
                scan_result=scan2["output"],
                error=f"Extracted ClamAV scan failed using {scan2.get('scanner')}",
            )
            return

        update_archive(con, archive_id, status="inventory")
        # A rerun creates a new /tmp/leak_processing/archive_<id>_<timestamp> path.
        # Remove old file inventory rows so the importer does not try deleted temp paths.
        con.execute("DELETE FROM leak_archive_files WHERE archive_id=?", (archive_id,))
        con.commit()
        files_found = inventory_files(con, archive_id, extract_dir)
        update_archive(con, archive_id, files_found=files_found)

        update_archive(con, archive_id, status="importing")
        archive = get_archive(con, archive_id)

        files = con.execute("""
            SELECT *
            FROM leak_archive_files
            WHERE archive_id=?
              AND status='queued'
              AND score >= 20
            ORDER BY score DESC, size_bytes DESC
        """, (archive_id,)).fetchall()

        total_imported = 0
        files_processed = 0

        for file_row in files:
            total_imported += process_file(con, archive, file_row)
            files_processed += 1

        update_archive(
            con,
            archive_id,
            status="processed",
            files_processed=files_processed,
            entities_imported=total_imported,
        )

    except Exception as e:
        update_archive(con, archive_id, status="error", error=str(e))

    finally:
        shutil.rmtree(work_dir, ignore_errors=True)
        con.close()

def resolve_anonfilesnew(url):
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Referer": url
    }

    r = requests.get(url, headers=headers, timeout=(10, 30))
    r.raise_for_status()

    soup = BeautifulSoup(r.text, "html.parser")

    hidden_url = None
    for a in soup.find_all("a", href=True):
        href = a.get("href", "")
        if "/cdn-cgi/content" in href:
            hidden_url = urljoin(url, href)
            break

    if not hidden_url:
        raise RuntimeError("Could not resolve anonfilesnew hidden cdn-cgi link")

    r2 = requests.get(hidden_url, headers=headers, allow_redirects=False, timeout=(10, 30))

    loc = r2.headers.get("Location") or r2.headers.get("location")
    if loc:
        final_url = urljoin(hidden_url, loc)
        if "/content/v4" in final_url:
            return final_url
        return final_url

    # fallback: sometimes content/v4 is inside response body
    m = re.search(r'https?://[^"\']+/content/v4\?s=[^"\'<>\s]+', r2.text)
    if m:
        return m.group(0)

    m = re.search(r'/content/v4\?s=[^"\'<>\s]+', r2.text)
    if m:
        return urljoin(hidden_url, m.group(0))

    raise RuntimeError("Could not resolve anonfilesnew content/v4 link")


HOST_RESOLVERS = {
    "anonfilesnew.com": resolve_anonfilesnew,
}

def resolve_download_url(url):
    parsed = urlparse(url)
    host = parsed.netloc.lower().replace("www.", "")
    resolver = HOST_RESOLVERS.get(host)
    if resolver:
        return resolver(url)
    return url

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--job-id", type=int, required=True)
    args = parser.parse_args()
    process_archive(args.job_id)


if __name__ == "__main__":
    main()
